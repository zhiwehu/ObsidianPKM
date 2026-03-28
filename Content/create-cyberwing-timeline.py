#!/usr/bin/env python3
"""
创建 Cyberwing 产品路线图 PPT - 表格时间线样式
参考用户提供的图片格式：按月份分列，按产品线/能力分行
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
# 16:9 幻灯片尺寸
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

# 颜色定义
def rgb(r, g, b):
    return RGBColor(r, g, b)

COLORS = {
    'header_bg': rgb(243, 244, 246),  # 浅灰色背景
    'border': rgb(209, 213, 219),
    'text_dark': rgb(31, 41, 55),
    'text_medium': rgb(107, 114, 128),
    'text_light': rgb(156, 163, 175),
    'accent_blue': rgb(59, 130, 246),
    'accent_orange': rgb(249, 115, 22),
    'accent_green': rgb(34, 197, 94),
    'row_bkg_1': rgb(255, 255, 255),
    'row_bkg_2': rgb(249, 250, 251),
}

# 月份定义（从 3 月到 8 月）
MONTHS = ['March', 'April', 'May', 'June', 'July', 'August']
MONTH_WIDTHS = [1.0, 2.0, 2.0, 1.8, 1.8, 1.8]  # 每个月份列的宽度

# 行定义
ROWS = [
    {'name': '技术能力', 'height': 1.8},
    {'name': '2C', 'height': 1.8},
    {'name': '2B', 'height': 1.5},
]

# 数据内容
DATA = {
    '技术能力': {
        'March': '',
        'April': '系统架构调整\n拾音系统优化\n声纹识别功能\n长短期记忆系统\n用户画像系统',
        'May': '声纹识别功能\n长短期记忆系统\n用户画像系统',
        'June': '',
        'July': '',
        'August': '',
    },
    '2C': {
        'March': '3/25 VIP 发货\n3/31 SS\nSW 1.0',
        'April': '4/10 后摩 LQ50 软硬件集成\n4/15 充值与积分管理\n4/30 文搜音频/视频能力\n4/30 HDMI 多媒体功能\nSW 1.1',
        'May': '5/10 NAS PC 支持 AI 搜索',
        'June': 'SW 2.0',
        'July': '',
        'August': '',
    },
    '2B': {
        'March': '摄影垂域智能相册 PRD',
        'April': '摄影垂域智能相册开发',
        'May': '5/29 摄影垂域智能相册发布\n自动备份/AI 自动整理/人脸识别\n/自然语言检索/AI 去闭眼/模糊/\n表情异常图片等',
        'June': '',
        'July': '',
        'August': '',
    },
}

# 起始位置
start_x = Inches(0.3)
start_y = Inches(0.5)
table_width = prs.slide_width - Inches(0.6)
row_heights = [1.8, 1.8, 1.5]  # 每行的高度

# 绘制标题
title = slide.shapes.add_textbox(start_x, Inches(0.15), Inches(8), Inches(0.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Cyberwing 产品路线图"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['text_dark']
p.font.name = 'Arial'

# 计算列位置
col_positions = [start_x]
for i, width in enumerate(MONTH_WIDTHS):
    col_positions.append(col_positions[-1] + Inches(width))

# 绘制表格
current_y = start_y

# 表头行（月份）
header_y = current_y
header_height = Inches(0.4)

# 绘制左上角空白单元格（行标题区域）
cell = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    start_x, header_y,
    Inches(1.5), header_height
)
cell.fill.solid()
cell.fill.fore_color.rgb = COLORS['header_bg']
cell.line.color.rgb = COLORS['border']
cell.line.width = Pt(1)

# 绘制月份表头
for i, month in enumerate(MONTHS):
    cell_x = start_x + Inches(1.5) + sum(Inches(w) for w in MONTH_WIDTHS[:i])
    cell_width = Inches(MONTH_WIDTHS[i])

    cell = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cell_x, header_y,
        cell_width, header_height
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLORS['header_bg']
    cell.line.color.rgb = COLORS['border']
    cell.line.width = Pt(1)

    # 添加月份文字
    textbox = slide.shapes.add_textbox(
        cell_x, header_y, cell_width, header_height
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = month
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_medium']
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

current_y += header_height

# 绘制数据行
for row_idx, row in enumerate(ROWS):
    row_height = Inches(row['height'])
    row_name = row['name']

    # 行标题单元格
    title_cell = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        start_x, current_y,
        Inches(1.5), row_height
    )
    title_cell.fill.solid()
    title_cell.fill.fore_color.rgb = COLORS['header_bg']
    title_cell.line.color.rgb = COLORS['border']
    title_cell.line.width = Pt(1)

    # 行标题文字
    title_textbox = slide.shapes.add_textbox(
        start_x, current_y, Inches(1.5), row_height
    )
    tf = title_textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = row_name
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_dark']
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.1)

    # 数据单元格
    for col_idx, month in enumerate(MONTHS):
        cell_x = start_x + Inches(1.5) + sum(Inches(w) for w in MONTH_WIDTHS[:col_idx])
        cell_width = Inches(MONTH_WIDTHS[col_idx])

        content = DATA[row_name].get(month, '')

        # 绘制单元格边框
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            cell_x, current_y,
            cell_width, row_height
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['row_bkg_1'] if row_idx % 2 == 0 else COLORS['row_bkg_2']
        cell.line.color.rgb = COLORS['border']
        cell.line.width = Pt(1)

        # 添加内容文字
        if content:
            textbox = slide.shapes.add_textbox(
                cell_x + Inches(0.08), current_y + Inches(0.08),
                cell_width - Inches(0.16), row_height - Inches(0.16)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            lines = content.split('\n')
            for line_idx, line in enumerate(lines):
                p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(9)
                p.font.color.rgb = COLORS['text_dark']
                p.font.name = 'Arial'
                if line.startswith('3/') or line.startswith('4/') or line.startswith('5/'):
                    p.font.bold = True
                p.space_after = Pt(2)

    current_y += row_height

# 添加底部说明
footer_y = current_y + Inches(0.2)
footer = slide.shapes.add_textbox(start_x, footer_y, Inches(6), Inches(0.3))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "TIMELINE 2026"
p.font.size = Pt(10)
p.font.color.rgb = COLORS['text_light']
p.font.name = 'Arial'

# 保存文件
output_path = '/Users/huzhiwei/Library/Mobile Documents/iCloud~md~obsidian/Documents/PKM/Content/Cyberwing-产品路线图 - 表格时间线.pptx'
prs.save(output_path)
print(f"✅ PPT 已保存到：{output_path}")
